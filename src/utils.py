
import os
from datetime import datetime
from pptx.dml.color import RGBColor
from pptx.text.text import TextFrame
from bs4 import BeautifulSoup
from bs4.element import NavigableString, Tag
import json
import yaml


ASSETS_DIR = 'assets'
PPT_TEMPLATE = 'Oracle_PPT-template_FY26_blank.pptx'
LAYOUT_TEMPLATE = 'layouts_template.yaml'
PROMPT_TEMPLATE = 'prompt_template.txt'


curdir = os.path.dirname(os.path.abspath(__file__))
assets_dir = os.path.join(curdir, ASSETS_DIR)
ppt_template = os.path.join(assets_dir,PPT_TEMPLATE)
layout_template = os.path.join(assets_dir,LAYOUT_TEMPLATE)
prompt_template = os.path.join(assets_dir,PROMPT_TEMPLATE)

def generate_timestamped_filename(
    save_location:str,
    name: str, 
    extension: str = "pptx") -> str:    
    timestamp = datetime.now() .strftime("%Y%m%d_%H%M%S")
    ext = extension.lstrip('.')
    filename = f"{name}_{timestamp}.{ext}"
    return os.path.join(save_location, filename)

def load_layout_template() -> dict:
    with open(layout_template, "r",encoding="utf-8") as f:
        yaml_data = f.read()
    layouts_template = yaml.safe_load(yaml_data)
    return layouts_template

def load_prompt_template() -> str:
    with open(prompt_template, "r",encoding="utf-8") as f:
        prompt_str = f.read()
    return prompt_str

def get_layouts() -> str:
    """
    Get the layouts template as text prompt.
    """
    layouts_template = load_layout_template()
    layout_dict = {}
    for k,v in layouts_template.items():
        placeholder_dict = {}        
        if v.get("description", "").strip() == "":
            continue
        layout_alias = v["alias"]
        if v.get("placeholders", []):
            for pk,pv in v["placeholders"].items():
                if pv.get("description", "").strip() == "":
                    continue
                alias = pv["alias"]
                placeholder_dict[alias] = pv["description"]
            layout_dict[layout_alias] = {
                    "description":v["description"],
                    "placeholders": placeholder_dict
                }
        else:
            layout_dict[layout_alias] = {
                    "description":v["description"],
                    "placeholders": "none"
                }
    layout_str = json.dumps(layout_dict, indent=4)
    return layout_str

def get_reverse_index() -> dict:
    layouts_template = load_layout_template()
    reverse_index = {}
    for k,v in layouts_template.items():
        alias = v["alias"]
        index = v["index"]
        reverse_p_index = {}
        if "placeholders" in v:
            for pk,pv in v["placeholders"].items():
                for ik,iy in pv["place_holder_index"].items():
                    p_alias = pv["alias"]
                    p_index = pv["place_holder_index"]
                    reverse_p_index[p_alias] = p_index
        reverse_index[alias] = {"layout_index":index, "placeholders": reverse_p_index}

    return reverse_index

def hex_to_rgb(
    hex_color: str
    ) -> tuple:
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2],16)
    g = int(hex_color[2:4],16)
    b = int(hex_color[4:6],16)
    return (r,g,b)

def parse_style(
    style_str: str
    ) -> dict:
    style_dict = {}
    items = style_str.strip(";").split(";")
    if items:
        for item in items:
            k,v = item.split(":")
            style_dict[k.strip()] = v.strip()
    return style_dict

def process_html(
    html: str, 
    p:TextFrame.paragraphs
    ) -> None:
    soup = BeautifulSoup(html, 'html.parser')
    p_tag = soup.find('p')
    s_tag = soup.find('span')
    if not p_tag and not s_tag:
        run = p.add_run()
        run.text = soup.get_text()
    else:
        if p_tag:
            if p_tag.has_attr('level'):
                p.level = int(p_tag.get('level'))                
            p_tag.unwrap()
        if s_tag:
            for element in soup.contents:
                run = p.add_run()
                run.text = element.get_text()
                if isinstance(element, NavigableString):
                    pass
                elif isinstance(element, Tag) and element.name == 'span':
                    if element.has_attr('data-link'):
                        run.hyperlink.address = element.get('data-link')
                    if element.has_attr('style'):
                        style_str = element.get('style', '')
                        style_dict = parse_style(style_str)
                        #print(style_dict)
                        if style_dict.get('font-weight') == 'bold':
                            run.font.bold = True
                            #print("bold")
                        if style_dict.get('font-style') == 'italic':
                            run.font.italic = True
                        if style_dict.get('text-decoration') == 'underline':
                            run.font.underline = True
                            #print("underline")
                        if style_dict.get('color'):
                            hex_color = style_dict.get('color')
                            run.font.color.rgb = RGBColor(*hex_to_rgb(hex_color))
                    #print(p.text)
        else:
            run = p.add_run()
            run.text = soup.get_text()
            