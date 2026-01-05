from pptx import Presentation

TEMPLATE_PATH = "./src/assets/Oracle_PPT-template_FY26_blank.pptx"

def print_layout_placeholders(target_layout_name=None):
    prs = Presentation(TEMPLATE_PATH)
    print("Total layouts:", len(prs.slide_layouts))
    for i, layout in enumerate(prs.slide_layouts):
        layout_name = layout.name
        print("="*60)
        print(f"Layout Index: {i}")
        print(f"Layout name: {layout_name}")
        if target_layout_name and (target_layout_name not in layout_name):
            continue
        for j, shape in enumerate(layout.placeholders):
            print(
                f"  Placeholder {j}:   "
                f"idx={shape.placeholder_format.idx}, "
                f"name='{shape.name}', "
                f"shape_id={shape.shape_id}, "
                f"type={shape.placeholder_format.type}"
            )
        print("-"*60)

if __name__ == "__main__":
    # Optionally pass a layout name substring to focus on one layout only
    # e.g., print_layout_placeholders("Title/Subtitle with 1 block")
    print_layout_placeholders()
