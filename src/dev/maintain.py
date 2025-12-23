from pptx import Presentation
import json
import yaml
import pandas as pd


def create_layout_all(filepath):
    prs = Presentation(filepath)
    layouts_list = []

    for i, layout in enumerate(prs.slide_layouts):    
        # if layout.name not in ["Light - Title_Pillar","Light - Divider","Light - Single metric","Light - Thank You","Light - Agenda Table","Light - Title 1 Column","Light - Title/Subtitle 1 Column"]:
        #     continue
        print(layout.name)
        if layout.name.lower().startswith("dark"):
            mode = "dark"
        elif layout.name.lower().startswith("light"):
            mode = "light"
        else:
            mode = "normal"

        shortname = layout.name
        for prefix in ["light","Light","LIGHT","Light","dark","Dark","DARK"]:            
            shortname = shortname.replace(prefix,"")
        shortname = shortname.strip(" ").strip("-").strip(" ")   
        layout_dict = {
                "index": i,
                "layout_name": layout.name,
                "mode": mode,
                "shortname": shortname,
                "description": "",
                "placeholders": []
        }
        
        for j,placeholder in enumerate(layout.placeholders):
            # if placeholder.name.startswith("Slide Number") or placeholder.name.startswith("Footer") or placeholder.name.startswith("Date"):
            #     continue
            layout_dict["placeholders"].append(
                {
                "index":j,
                "place_holder_index":placeholder.placeholder_format.idx,
                "placeholder_name": placeholder.name,
                "placeholder_id": placeholder.shape_id,
                "alias": placeholder.name,
                "description": ""
                }
            )
        
        layouts_list.append(layout_dict)

    yaml_data = yaml.dump({"layouts": layouts_list}, sort_keys=False)
    with open("layouts_all.yaml", "w") as f:
        f.write(yaml_data)

def create_layout_description():
    with open("layouts_all.yaml", "r") as f:
        yaml_data = f.read()
    layouts = yaml.safe_load(yaml_data)["layouts"]

    layout_description = {}
    data = []

    for layout in layouts:
        shortname = layout["shortname"]
        mode = layout["mode"]
        if shortname not in layout_description.keys():
            layout_description[shortname] = {
                "disabled": True,
                "layout_name": {mode:layout["layout_name"]},
                "index": {mode: layout["index"]},
                "alias": shortname,
                "description": layout["description"]
            }

            placeholders = {}
            for placeholder in layout["placeholders"]:            
                placeholders[placeholder["index"]] = {
                    "disabled": True,
                    "name":placeholder["placeholder_name"],
                    "alias":placeholder["alias"],
                    "description":placeholder["description"],
                    "place_holder_index":{mode:placeholder["place_holder_index"]}
                    }

                
            layout_description[shortname]["placeholders"] = placeholders
        else:   
            layout_description[shortname]["index"][mode] = layout["index"]
            for placeholder in layout["placeholders"]:            
                try:
                    layout_description[shortname]["layout_name"][mode] = layout["layout_name"]
                    layout_description[shortname]["placeholders"][placeholder["index"]]["place_holder_index"][mode] = placeholder["place_holder_index"]

                    # data.append({
                    #     "shortname": shortname,
                    #     "layout_name": layout["layout_name"],
                    #     "index": layout["index"],
                    #     "alias": shortname,
                    #     "description": layout["description"],
                    #     "placeholders_name": placeholder["placeholder_name"],
                    #     "placeholders_alias": placeholder["alias"],
                    #     "placeholders_description": placeholder["description"],
                    #     "placeholders_place_holder_index": placeholder["place_holder_index"]
                    # })
                except Exception as e:
                    print(e)
                    print(layout_description[shortname])
                    print(layout)
                    break
        

    with open("layouts_description.yaml", "w") as f:
        f.write(yaml.dump(layout_description, sort_keys=False))
    
    # df = pd.DataFrame(data)
    # df.to_excel("layouts_description.xlsx", index=  True)


def update_layout_description():
    with open("layouts_description_edit.yaml", "r") as f:
        yaml_data = f.read()
    template = yaml.safe_load(yaml_data)

    with open("layouts_description.yaml", "r") as f:
        yaml_data = f.read()
    layouts = yaml.safe_load(yaml_data)

    for k,v in layouts.items():
        if k in template.keys():
            layouts[k]["alias"] = template[k]["alias"]
            layouts[k]["description"] = template[k]["description"]
            layouts[k]["placeholders"] = template[k]["placeholders"]
        else:
            layouts[k]["disabled"] = True

    with open("layouts_description.yaml", "w") as f:
        f.write(yaml.dump(layouts, sort_keys=False))
