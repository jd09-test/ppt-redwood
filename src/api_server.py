from dotenv import load_dotenv
load_dotenv()

from fastapi import FastAPI, UploadFile, File, HTTPException
from pydantic import BaseModel
from typing import List, Optional
import os
import json
import tempfile

from src.utils import (
    get_layouts,
    load_prompt_template,
    generate_timestamped_filename,
    ppt_template,
    process_html,
    get_reverse_index,
)
from src.model import PresentationContent
from pptx import Presentation

import openai
import cloudinary
import cloudinary.uploader
from docx import Document
import edge_tts
import asyncio

cloudinary.config(
    cloud_name=os.environ.get("CLOUDINARY_CLOUD_NAME"),
    api_key=os.environ.get("CLOUDINARY_API_KEY"),
    api_secret=os.environ.get("CLOUDINARY_API_SECRET"),
    secure=True
)

app = FastAPI(
    title="PPT Generator API",
    description="HTTP API for generating PowerPoint presentations, Word documents, and adding cloud AI narration.",
    version="0.1.0"
)

class GeneratePresentationRequest(BaseModel):
    json_content: dict

# 1. get_presentation_rules
@app.get("/get_presentation_rules")
async def get_presentation_rules():
    try:
        layouts_str = get_layouts()
        prompt_str = load_prompt_template()
        from string import Template
        prompt_template = Template(prompt_str)
        rule_prompt = prompt_template.substitute(layouts_description=layouts_str)
        return {"rule_prompt": rule_prompt}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating rules: {str(e)}")

# Re-add this for /generate_json_content
class GenerateJsonContentRequest(BaseModel):
    user_text: str
    theme_mode: Optional[str] = "light"

@app.post("/generate_json_content")
async def generate_json_content(data: GenerateJsonContentRequest):
    openai.api_key = os.environ.get("OPENAI_API_KEY")
    if not openai.api_key:
        raise HTTPException(status_code=503, detail="OpenAI API key not set in environment variable OPENAI_API_KEY.")
    try:
        layouts_str = get_layouts()
        prompt_str = load_prompt_template()
        from string import Template
        prompt_template = Template(prompt_str)
        system_prompt = prompt_template.substitute(layouts_description=layouts_str)
        user_prompt = (
            f"Input for presentation: {data.user_text}\n"
            f"Please use theme_mode: '{data.theme_mode}'."
        )
        completion = openai.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.3,
            max_tokens=4096
        )
        ai_json = completion.choices[0].message.content
        import re
        pattern = r"```(?:json)?\s*([\s\S]+?)\s*```"
        match = re.search(pattern, ai_json)
        if match:
            json_text = match.group(1)
        else:
            json_text = ai_json
        result = json.loads(json_text)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"OpenAI completion error: {str(e)}")

# 2. generate_presentation API
@app.post("/generate_presentation")
async def generate_presentation(data: GeneratePresentationRequest):
    try:
        prs_content = PresentationContent(**data.json_content)
        slides_content = prs_content.slides
        theme_mode = prs_content.theme_mode

        message = {"slides_count": 0}
        prs = Presentation(ppt_template)
        reverse_index = get_reverse_index()
        import re
        layout_keys = set(reverse_index.keys())
        def normalize_layout_key(in_key):
            if in_key in layout_keys:
                return in_key
            alt = in_key.replace("_", "/").strip()
            if alt in layout_keys:
                return alt
            for key in layout_keys:
                if key.replace(" ", "").replace("/", "").lower() == in_key.replace(" ", "").replace("_", "").lower():
                    return key
            return None

        for slide_content in slides_content:
            orig_layout = slide_content.layout
            layout_name = normalize_layout_key(orig_layout)
            if not layout_name:
                raise HTTPException(
                    status_code=422,
                    detail=f"Layout '{orig_layout}' not recognized. Valid layouts: {sorted(layout_keys)}"
                )
            layout_index = reverse_index[layout_name]["layout_index"][theme_mode]
            placeholders_reverse = reverse_index[layout_name]["placeholders"]
            slide_layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(slide_layout)
            for placeholder in slide_content.placeholders:
                name = placeholder.placeholder_name
                content = placeholder.content
                try:
                    index = placeholders_reverse[name][theme_mode]
                    ph = slide.placeholders.__getitem__(idx=index)
                    text_frame = ph.text_frame
                    text_frame.clear()
                    text_frame._element.remove(text_frame.paragraphs[0]._p)
                    if isinstance(content, list):
                        for c in content:
                            p = text_frame.add_paragraph()
                            process_html(c, p)
                except Exception as e:
                    break
            # Add speaker notes if provided
            if hasattr(slide_content, "speaker_notes") and slide_content.speaker_notes:
                slide.notes_slide.notes_text_frame.text = slide_content.speaker_notes
            message["slides_count"] += 1

        # Save PPTX to in-memory buffer, upload to Cloudinary, do not store to disk
        import io
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        try:
            upload_result = cloudinary.uploader.upload(
                pptx_buffer,
                resource_type="raw",
                public_id=prs_content.filename,
                folder="presentations"
            )
            pptx_url = upload_result.get("secure_url")
        except Exception as e:
            pptx_url = None

        return { "cloudinary_url": pptx_url }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating PPTX: {str(e)}")

# 3. generate_word_doc API remains
class GenerateWordDocRequest(BaseModel):
    content: str
    filename: Optional[str] = "blog_post"

@app.post("/generate_word_doc")
async def generate_word_doc(data: GenerateWordDocRequest):
    try:
        doc = Document()
        paragraphs = data.content.split("\n\n")
        for para in paragraphs:
            doc.add_paragraph(para.strip())
        import io
        docx_buffer = io.BytesIO()
        doc.save(docx_buffer)
        docx_buffer.seek(0)
        try:
            upload_result = cloudinary.uploader.upload(
                docx_buffer,
                resource_type="raw",
                public_id=data.filename,
                folder="word_docs"
            )
            cloud_url = upload_result.get("secure_url")
        except Exception as upload_err:
            raise HTTPException(status_code=500, detail=f"Cloudinary upload error: {str(upload_err)}")
        return { "cloudinary_url": cloud_url }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating DOCX: {str(e)}")

# 4. generate_ppt_with_audio unchanged

@app.post("/generate_ppt_with_audio")
async def generate_ppt_with_audio(data: GeneratePresentationRequest):
    try:
        # Now support: "audio" inside json_content root (not at the outer level)
        # Default: "female" if not set
        audio_gender = data.json_content.get("audio", "female")
        prs_content = PresentationContent(**{k: v for k, v in data.json_content.items() if k not in ["audio"]})
        slides_content = prs_content.slides
        theme_mode = prs_content.theme_mode

        if audio_gender == "male":
            voice = "en-US-AndrewMultilingualNeural"
        else:
            voice = "en-IN-NeerjaNeural"

        message = {"slides_count": 0}
        prs = Presentation(ppt_template)
        reverse_index = get_reverse_index()

        import re
        layout_keys = set(reverse_index.keys())
        def normalize_layout_key(in_key):
            if in_key in layout_keys:
                return in_key
            alt = in_key.replace("_", "/").strip()
            if alt in layout_keys:
                return alt
            for key in layout_keys:
                if key.replace(" ", "").replace("/", "").lower() == in_key.replace(" ", "").replace("_", "").lower():
                    return key
            return None

        import zipfile

        slide_audio_urls = []
        slide_audio_files = []
        for i, slide_content in enumerate(slides_content, start=1):
            orig_layout = slide_content.layout
            layout_name = normalize_layout_key(orig_layout)
            if not layout_name:
                raise HTTPException(
                    status_code=422,
                    detail=f"Layout '{orig_layout}' not recognized. Valid layouts: {sorted(layout_keys)}"
                )
            layout_index = reverse_index[layout_name]["layout_index"][theme_mode]
            placeholders_reverse = reverse_index[layout_name]["placeholders"]
            slide_layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(slide_layout)
            for placeholder in slide_content.placeholders:
                name = placeholder.placeholder_name
                content = placeholder.content
                try:
                    index = placeholders_reverse[name][theme_mode]
                    ph = slide.placeholders.__getitem__(idx=index)
                    text_frame = ph.text_frame
                    text_frame.clear()
                    text_frame._element.remove(text_frame.paragraphs[0]._p)
                    if isinstance(content, list):
                        for c in content:
                            p = text_frame.add_paragraph()
                            process_html(c, p)
                except Exception as e:
                    break
            # Per-slide audio logic
            slide_audio_url = None
            slide_audio_tmpfile = None
            if hasattr(slide_content, "speaker_notes") and slide_content.speaker_notes:
                slide.notes_slide.notes_text_frame.text = slide_content.speaker_notes
                # Generate audio for this slide's speaker notes only
                try:
                    with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as tmpfile:
                        tmpfilename = tmpfile.name
                    communicate = edge_tts.Communicate(slide_content.speaker_notes, voice=voice)
                    await communicate.save(tmpfilename)
                    with open(tmpfilename, "rb") as f:
                        upload_result = cloudinary.uploader.upload(
                            f,
                            resource_type="video",
                            public_id=f"{prs_content.filename}_slide{i}_audio",
                            folder="ppt_audio"
                        )
                        slide_audio_url = upload_result.get("secure_url")
                    slide_audio_tmpfile = tmpfilename
                except Exception as audio_err:
                    slide_audio_url = None
                    slide_audio_tmpfile = None
            slide_audio_urls.append(slide_audio_url)
            slide_audio_files.append(slide_audio_tmpfile)
            message["slides_count"] += 1

        # Create ZIP of all non-None slide audio files
        audio_zip_url = None
        audio_zip_tmpfile = None
        try:
            zip_path = tempfile.NamedTemporaryFile(suffix=".zip", delete=False).name
            with zipfile.ZipFile(zip_path, "w") as zipf:
                for j, fpath in enumerate(slide_audio_files, start=1):
                    if fpath and os.path.exists(fpath):
                        # Name inside ZIP: slide{j}.mp3
                        zipf.write(fpath, f"slide{j}.mp3")
            with open(zip_path, "rb") as zf:
                upload_result = cloudinary.uploader.upload(
                    zf,
                    resource_type="raw",
                    public_id=f"{prs_content.filename}_audio_bundle",
                    folder="ppt_audio",
                    use_filename=True,
                    unique_filename=False,
                    overwrite=True,
                    type="upload"  # ensure public accessibility
                )
                base_url = upload_result.get("secure_url")
                audio_zip_url = base_url + "?fl_attachment"
            audio_zip_tmpfile = zip_path
        except Exception as zip_err:
            audio_zip_url = None

        # Clean up temp mp3 files and zip
        for f in slide_audio_files:
            if f and os.path.exists(f):
                try:
                    os.remove(f)
                except Exception:
                    pass
        if audio_zip_tmpfile and os.path.exists(audio_zip_tmpfile):
            try:
                os.remove(audio_zip_tmpfile)
            except Exception:
                pass

        # Remove single audio file logic

        import io
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        pptx_url = None
        try:
            upload_result = cloudinary.uploader.upload(
                pptx_buffer,
                resource_type="raw",
                public_id=prs_content.filename,
                folder="presentations"
            )
            pptx_url = upload_result.get("secure_url")
        except Exception as cloud_err:
            pptx_url = None

        return {
            "cloudinary_url": pptx_url,
            "slide_audio_urls": slide_audio_urls,
            "audio_zip_url": audio_zip_url
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating PPTX/audio: {str(e)}")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("src.api_server:app", host="0.0.0.0", port=int(os.environ.get("PORT", 8000)), reload=True)
